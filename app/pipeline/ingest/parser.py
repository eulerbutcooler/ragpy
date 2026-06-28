import logging
from collections import defaultdict
from pathlib import Path

from llama_index.core import Document

logger = logging.getLogger(__name__)

def _sanitize_text(text: str) -> str:
    # Convert surrogate pairs to actual codepoints and drop isolated surrogates
    # so downstream utf-8 encoding (e.g., llama_index hashing) never blows up.
    return text.encode("utf-16", "surrogatepass").decode("utf-16", "ignore")


def _sanitize_document(doc: Document) -> Document:
    """Reconstruct a Document with sanitized text.

    llama_index Document is a Pydantic model where ``text`` is a read-only
    property backed by ``text_resource``.  Mutating ``doc.text`` directly
    raises ``AttributeError: property 'text' of 'Document' object has no
    setter``.  Reconstructing preserves all important fields while swapping
    in clean text.
    """
    return Document(
        id_=doc.id_,
        text=_sanitize_text(doc.text),
        metadata=doc.metadata,
        embedding=doc.embedding,
        excluded_embed_metadata_keys=doc.excluded_embed_metadata_keys,
        excluded_llm_metadata_keys=doc.excluded_llm_metadata_keys,
        relationships=doc.relationships,
    )


def _table_to_markdown(raw_rows: list[list[str]]) -> str:
    if not raw_rows:
        return ""

    def dedup_row(row: list[str]) -> list[str]:
        if not row:
            return row
        out = [row[0]]
        for cell in row[1:]:
            if cell != out[-1]:
                out.append(cell)
        return out

    deduped = [dedup_row(r) for r in raw_rows]
    col_count = max(len(r) for r in deduped)

    def pad(row: list[str]) -> list[str]:
        return (row + [""] * col_count)[:col_count]

    header = pad(deduped[0])
    lines = [
        "| " + " | ".join(header) + " |",
        "|" + "|".join(["---"] * col_count) + "|",
    ]
    for row in deduped[1:]:
        lines.append("| " + " | ".join(
            c.strip().replace("\n", " ") for c in pad(row)
        ) + " |")
    return "\n".join(lines)

def parse_docx(file_path: Path) -> list[Document]:
    from docx import Document as DocxDocument
    from docx.table import Table

    doc = DocxDocument(str(file_path))
    parts: list[str] = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1]

        if tag == "p":
            para = next(
                (p for p in doc.paragraphs if p._element is element), None
            )
            if para and para.text.strip():
                parts.append(para.text.strip())

        elif tag == "tbl":
            table = Table(element, doc)
            raw_rows: list[list[str]] = []
            for row in table.rows:
                raw_rows.append([
                    cell.text.strip().replace("\n", " ")
                    for cell in row.cells
                ])

            if raw_rows:
                parts.append(_table_to_markdown(raw_rows))

    if not parts:
        logger.warning("docx_empty file=%s", file_path.name)
        return []

    text = _sanitize_text("\n\n".join(parts))
    logger.info("docx_parsed file=%s len=%d", file_path.name, len(text))
    return [Document(text=text)]

def _extract_slide_text(slide) -> str:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    parts: list[str] = []

    def process_shape(shape) -> None:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                process_shape(s)
            return
        if shape.has_table:
            raw_rows = [
                [
                    cell.text_frame.text.strip().replace("\n", " ")
                    for cell in row.cells
                ]
                for row in shape.table.rows
            ]
            if raw_rows:
                parts.append(_table_to_markdown(raw_rows))
            return
        if shape.has_text_frame:
            lines = []
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    indent = "  " * para.level if para.level else ""
                    lines.append(f"{indent}{line}")
            if lines:
                parts.append("\n".join(lines))

    for shape in slide.shapes:
        process_shape(shape)

    return "\n\n".join(parts)

def _extract_slide_notes(slide) -> str:
    try:
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf:
                return tf.text.strip()
    except Exception:
        pass
    return ""

def parse_pptx(file_path: Path) -> list[Document]:
    from pptx import Presentation

    prs = Presentation(str(file_path))
    total_slides = len(prs.slides)

    if total_slides == 0:
        logger.warning("pptx_empty file=%s", file_path.name)
        return []

    documents: list[Document] = []

    for i, slide in enumerate(prs.slides, 1):
        content = _extract_slide_text(slide)
        notes = _extract_slide_notes(slide)

        if not content.strip():
            logger.warning(
                "pptx_slide_no_text file=%s slide=%d",
                file_path.name, i,
            )
            continue

        full_text = _sanitize_text(f"Slide {i}:\n\n{content}")
        if notes:
            full_text += _sanitize_text(f"\n\nSpeaker notes:\n{notes}")

        documents.append(Document(
            text=full_text,
            metadata={"slide_number": i},
        ))

    logger.info(
        "pptx_parsed file=%s slides_indexed=%d/%d",
        file_path.name, len(documents), total_slides,
    )
    return documents

def parse_pdf(file_path: Path) -> list[Document]:
    from llama_index.readers.file import PyMuPDFReader

    documents = PyMuPDFReader().load_data(file_path=file_path)
    sanitized = [_sanitize_document(doc) for doc in documents]
    logger.info(
        "pdf_parsed file=%s pages=%d", file_path.name, len(sanitized)
    )
    return sanitized

_PARSERS = {
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".pdf":  parse_pdf,
}

_SUPPORTED_EXTENSIONS = set(_PARSERS.keys())

def parse_file(file_path: Path) -> list[Document]:
    ext = file_path.suffix.lower()
    parser = _PARSERS.get(ext)
    if parser is None:
        raise ValueError(f"Unsupported file type: {ext!r}")
    return parser(file_path)
